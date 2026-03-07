#!/usr/bin/env python3
"""
創建重慶美食PPT
使用python-pptx庫
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

def create_chongqing_food_ppt():
    """創建重慶美食PPT"""
    
    # 創建演示文稿
    prs = Presentation()
    
    # 設置幻燈片尺寸 (16:9)
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)
    
    # ========== 第一頁：標題頁 ==========
    slide_layout = prs.slide_layouts[0]  # 標題幻燈片
    slide1 = prs.slides.add_slide(slide_layout)
    
    # 設置背景顏色
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(178, 34, 34)  # 火鍋紅
    
    # 標題
    title = slide1.shapes.title
    title.text = "重慶美食之旅"
    title.text_frame.paragraphs[0].font.size = Pt(60)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # 白色
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 副標題
    subtitle = slide1.placeholders[1]
    subtitle.text = "麻辣之都的美食探索"
    subtitle.text_frame.paragraphs[0].font.size = Pt(32)
    subtitle.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 224)  # 淺黃色
    subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # 作者信息
    left = Inches(1)
    top = Inches(6.5)
    width = Inches(11.33)
    height = Inches(0.5)
    
    textbox = slide1.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    p = tf.add_paragraph()
    p.text = "Gordon Lui"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.RIGHT
    
    # ========== 第二頁：重慶美食簡介 ==========
    slide_layout = prs.slide_layouts[1]  # 標題和內容
    slide2 = prs.slides.add_slide(slide_layout)
    
    # 設置背景
    background = slide2.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 250, 240)  # 米白色
    
    # 標題
    title = slide2.shapes.title
    title.text = "重慶美食簡介"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(178, 34, 34)  # 火鍋紅
    title.text_frame.paragraphs[0].font.bold = True
    
    # 內容
    content = slide2.placeholders[1]
    tf = content.text_frame
    tf.clear()  # 清除默認文本
    
    # 第一點
    p = tf.add_paragraph()
    p.text = "📍 重慶：中國麻辣美食之都"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.font.bold = True
    p.level = 0
    
    # 第二點
    p = tf.add_paragraph()
    p.text = "🌍 地理位置：長江與嘉陵江交匯處"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.level = 1
    
    # 第三點
    p = tf.add_paragraph()
    p.text = "🌶️ 美食特色：麻辣鮮香，層次豐富"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.level = 1
    
    # 第四點
    p = tf.add_paragraph()
    p.text = "🏮 文化背景：碼頭文化與山城特色"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.level = 1
    
    # 第五點
    p = tf.add_paragraph()
    p.text = "🎯 代表美食：火鍋、小麵、酸辣粉"
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.level = 1
    
    # 添加裝飾形狀
    left = Inches(9)
    top = Inches(2)
    width = Inches(3)
    height = Inches(3)
    
    shape = slide2.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(220, 20, 60)  # 深紅色
    shape.line.color.rgb = RGBColor(255, 255, 255)
    shape.line.width = Pt(2)
    
    # 形狀內文字
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = "麻辣\n之都"
    p.font.size = Pt(32)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # ========== 第三頁：三大經典美食 ==========
    slide_layout = prs.slide_layouts[1]
    slide3 = prs.slides.add_slide(slide_layout)
    
    # 設置背景
    background = slide3.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 245, 238)  # 杏仁白
    
    # 標題
    title = slide3.shapes.title
    title.text = "三大經典美食"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(178, 34, 34)
    title.text_frame.paragraphs[0].font.bold = True
    
    # 內容 - 使用表格形式
    content = slide3.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # 火鍋部分
    p = tf.add_paragraph()
    p.text = "🥘 重慶火鍋"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(220, 20, 60)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• 起源：江邊碼頭工人的簡易飲食"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 特色：九宮格，牛油鍋底，麻辣鮮香"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 必點：毛肚、鴨腸、黃喉、腦花"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.level = 1
    
    # 小麵部分
    p = tf.add_paragraph()
    p.text = "🍜 重慶小麵"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(220, 20, 60)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• 特點：麻辣勁道，湯汁濃郁"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 種類：豌雜麵、牛肉麵、肥腸麵"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 精髓：辣椒油和花椒完美結合"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.level = 1
    
    # 酸辣粉部分
    p = tf.add_paragraph()
    p.text = "🌶️ 酸辣粉"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(220, 20, 60)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• 特色：酸辣開胃，粉條爽滑"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 配料：花生、榨菜、肉末、香菜"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 地位：街頭小吃代表"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.level = 1
    
    # 添加裝飾線
    left = Inches(0.5)
    top = Inches(5.5)
    width = Inches(12.33)
    height = Inches(0.1)
    
    line = slide3.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(178, 34, 34)
    line.line.color.rgb = RGBColor(178, 34, 34)
    
    # ========== 第四頁：美食地圖與推薦 ==========
    slide_layout = prs.slide_layouts[1]
    slide4 = prs.slides.add_slide(slide_layout)
    
    # 設置背景
    background = slide4.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 255, 240)  # 蜜瓜綠
    
    # 標題
    title = slide4.shapes.title
    title.text = "美食地圖與推薦"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(178, 34, 34)
    title.text_frame.paragraphs[0].font.bold = True
    
    # 內容
    content = slide4.placeholders[1]
    tf = content.text_frame
    tf.clear()
    
    # 必去美食街區
    p = tf.add_paragraph()
    p.text = "🗺️ 必去美食街區"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0, 100, 0)  # 深綠色
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• 解放碑商圈：傳統與現代美食交匯"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 洪崖洞：夜景與美食完美結合"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 磁器口古鎮：傳統小吃聚集地"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.level = 1
    
    # 餐廳推薦
    p = tf.add_paragraph()
    p.text = "🏆 餐廳推薦"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0, 100, 0)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "1. 珮姐老火鍋 - 網紅火鍋店，常年排隊"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "2. 李串串老火鍋 - 本地人最愛"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "3. 花市豌雜麵 - 重慶小麵代表"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "4. 好又來酸辣粉 - 經典酸辣粉"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.level = 1
    
    # 旅行小貼士
    p = tf.add_paragraph()
    p.text = "💡 旅行小貼士"
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0, 100, 0)
    p.font.bold = True
    p.level = 0
    
    p = tf.add_paragraph()
    p.text = "• 最佳季節：秋季（9-11月）"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 辣度選擇：可要求「微辣」或「少辣」"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 搭配飲品：唯怡豆奶、王老吉"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• 必備物品：腸胃藥、濕紙巾"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(50, 50, 50)
    p.level = 1
    
    # 添加結尾裝飾
    left = Inches(10)
    top = Inches(6)
    width = Inches(2)
    height = Inches(1)
    
    shape = slide4.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(178, 34, 34)
    shape.line.color.rgb = RGBColor(255, 255, 255)
    shape.line.width = Pt(2)
    
    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = "享受\n美食！"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # ========== 保存文件 ==========
    output_path = "/Users/gordonlui/.openclaw/workspace/重慶美食之旅.pptx"
    prs.save(output_path)
    
    print(f"✅ PPT已創建: {output_path}")
    print(f"📊 總頁數: {len(prs.slides)} 頁")
    
    return output_path

def main():
    """主函數"""
    print("=" * 70)
    print("🎯 創建重慶美食PPT")
    print("=" * 70)
    
    try:
        # 檢查python-pptx是否安裝
        from pptx import Presentation
        print("✅ python-pptx庫已可用")
        
        # 創建PPT
        ppt_path = create_chongqing_food_ppt()
        
        # 顯示文件信息
        file_size = os.path.getsize(ppt_path) / 1024  # KB
        print(f"📁 文件大小: {file_size:.1f} KB")
        
        print(f"\n📋 PPT內容概要:")
        print(f"  第1頁: 標題頁 - 重慶美食之旅")
        print(f"  第2頁: 重慶美食簡介")
        print(f"  第3頁: 三大經典美食")
        print(f"  第4頁: 美食地圖與推薦")
        
        print(f"\n🎨 設計特色:")
        print(f"  • 主題顏色: 火鍋紅搭配暖色系")
        print(f"  • 字體大小: 層次分明，易於閱讀")
        print(f"  • 圖標使用: 美食相關emoji圖標")
        print(f"  • 佈局設計: 簡潔專業，重點突出")
        
        return ppt_path
        
    except ImportError:
        print("❌ 錯誤: 需要安裝python-pptx庫")
        print("💡 安裝命令: pip install python-pptx")
        return None
    except Exception as e:
        print(f"❌ 創建PPT時出錯: {e}")
        return None

if __name__ == "__main__":
    ppt_path = main()
    
    if ppt_path:
        print(f"\n✅ 重慶美食PPT創建成功！")
        print(f"📂 文件位置: {ppt_path}")
        print(f"\n📧 下一步: 發送到 gordonlct125@gmail.com")
        print("=" * 70)
    else:
        print(f"\n❌ PPT創建失敗")
        print("=" * 70)